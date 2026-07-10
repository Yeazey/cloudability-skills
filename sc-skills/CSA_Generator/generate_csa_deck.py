#!/usr/bin/env python3
"""
CSA ROI Assessment PowerPoint Generator
Universal script for generating CSA ROI Assessment decks from Excel files.

Usage:
    python3 generate_csa_deck.py <path_to_excel_file>

Requirements:
    - CSA ROI Assessment Template.pptx must be in the same directory as this script
    - Excel file must have 'Executive Summary' sheet with standard CSA format
"""

from pptx import Presentation
from pptx.util import Pt
import pandas as pd
import sys
import os
import re

def find_row_by_label(df, label_text, col=1):
    """Find row index by searching for label text in specified column"""
    for idx in range(len(df)):
        cell_value = df.iloc[idx, col]
        if pd.notna(cell_value) and label_text in str(cell_value):
            return idx
    return None

def safe_float(value):
    """Safely convert value to float, return 0 if NaN or invalid"""
    if pd.isna(value):
        return 0.0
    try:
        return float(value)
    except (ValueError, TypeError):
        return 0.0

def set_table_font_size(table, font_size):
    """Set font size for all cells in a table"""
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(font_size)

def extract_customer_name(excel_path):
    """Extract customer name from Excel filename"""
    filename = os.path.basename(excel_path)
    
    # Pattern 1: "Financial Planning Report - Account# - Company Name - ..."
    match = re.match(r'^Financial Planning Report\s*-\s*\d+\s*-\s*(.+?)\s*-\s*\d+', filename)
    if match:
        return match.group(1).strip()
    
    # Pattern 2: "Company Name - Account# - ..."
    match = re.match(r'^(.+?)\s*-\s*\d+', filename)
    if match:
        return match.group(1).strip()
    
    return "Customer"

def main():
    if len(sys.argv) < 2:
        print("Usage: python3 generate_csa_deck.py <path_to_excel_file>")
        sys.exit(1)
    
    excel_file = sys.argv[1]
    
    if not os.path.exists(excel_file):
        print(f"ERROR: Excel file not found: {excel_file}")
        sys.exit(1)
    
    # Find template
    script_dir = os.path.dirname(os.path.abspath(__file__))
    template_pptx = os.path.join(script_dir, "CSA ROI Assessment Template.pptx")
    
    if not os.path.exists(template_pptx):
        print(f"ERROR: Template not found: {template_pptx}")
        print("Please ensure 'CSA ROI Assessment Template.pptx' is in the same directory as this script.")
        sys.exit(1)
    
    # Extract customer name and create output filename
    customer_name = extract_customer_name(excel_file)
    excel_dir = os.path.dirname(excel_file)
    output_pptx = os.path.join(excel_dir, f"{customer_name} CSA ROI Assessment.pptx")
    
    try:
        print("="*80)
        print("CSA ROI ASSESSMENT GENERATOR")
        print("="*80)
        print(f"\nInput Excel: {excel_file}")
        print(f"Customer: {customer_name}")
        print(f"Output: {output_pptx}")
        
        # Load Executive Summary
        print("\n1. Loading Executive Summary...")
        exec_summary = pd.read_excel(excel_file, sheet_name='Executive Summary', header=None)
        
        # SLIDE 9 DATA EXTRACTION
        print("2. Extracting Slide 9 data...")
        
        edp_row = find_row_by_label(exec_summary, "EDP (On-demand usage)")
        current_coverage_row = find_row_by_label(exec_summary, "Current Coverage")
        flexibility_as_is_row = find_row_by_label(exec_summary, "Commitment cash flow flexibility - As-Is")
        flexibility_row = find_row_by_label(exec_summary, "Commitment cash flow flexibility")
        current_state_row = find_row_by_label(exec_summary, "Current state - savings")
        cloudwiry_savings_row = find_row_by_label(exec_summary, "Cloudwiry - savings")
        projected_discount_row = find_row_by_label(exec_summary, "Projected Net Effective Discount")
        stable_row = find_row_by_label(exec_summary, "stable", col=10)
        fluctuations_row = find_row_by_label(exec_summary, "fluctuations", col=10)
        stable_equiv_row = find_row_by_label(exec_summary, "Stable OD Equivalent", col=10)
        
        # Extract values
        current_edp = safe_float(exec_summary.iloc[edp_row, 2]) if edp_row else 0.0
        current_coverage = safe_float(exec_summary.iloc[current_coverage_row, 2]) if current_coverage_row else 0.0
        current_flexibility = safe_float(exec_summary.iloc[flexibility_as_is_row, 3]) if flexibility_as_is_row else 0.0
        current_savings = safe_float(exec_summary.iloc[current_state_row, 3]) if current_state_row else 0.0
        
        projected_edp = current_edp
        projected_coverage = safe_float(exec_summary.iloc[projected_discount_row, 2]) if projected_discount_row else 0.0
        
        # Find non-As-Is flexibility row
        if flexibility_as_is_row is not None and flexibility_row is not None:
            if flexibility_row == flexibility_as_is_row:
                for idx in range(flexibility_as_is_row + 1, len(exec_summary)):
                    cell_value = exec_summary.iloc[idx, 1]
                    if pd.notna(cell_value) and "Commitment cash flow flexibility" in str(cell_value) and "As-Is" not in str(cell_value):
                        flexibility_row = idx
                        break
        
        projected_flexibility = safe_float(exec_summary.iloc[flexibility_row, 3]) if flexibility_row else 0.0
        projected_savings = safe_float(exec_summary.iloc[cloudwiry_savings_row, 3]) if cloudwiry_savings_row else 0.0
        
        stable_pct = safe_float(exec_summary.iloc[stable_row, 11]) if stable_row else 0.0
        fluctuations_pct = safe_float(exec_summary.iloc[fluctuations_row, 11]) if fluctuations_row else 0.0
        stable_equivalent = safe_float(exec_summary.iloc[stable_equiv_row, 11]) if stable_equiv_row else 0.0
        
        # SLIDE 11 DATA EXTRACTION
        print("3. Extracting Slide 11 data...")
        
        incremental_row = find_row_by_label(exec_summary, "Incremental")
        pct_increase_row = find_row_by_label(exec_summary, "% Increase")
        flexibility_final_row = find_row_by_label(exec_summary, "Flexibility")
        
        current_year1 = safe_float(exec_summary.iloc[current_state_row, 3]) if current_state_row else 0.0
        current_year2 = safe_float(exec_summary.iloc[current_state_row, 4]) if current_state_row else 0.0
        current_year3 = safe_float(exec_summary.iloc[current_state_row, 5]) if current_state_row else 0.0
        
        csa_year1 = safe_float(exec_summary.iloc[cloudwiry_savings_row, 3]) if cloudwiry_savings_row else 0.0
        csa_year2 = safe_float(exec_summary.iloc[cloudwiry_savings_row, 4]) if cloudwiry_savings_row else 0.0
        csa_year3 = safe_float(exec_summary.iloc[cloudwiry_savings_row, 5]) if cloudwiry_savings_row else 0.0
        
        incr_year1 = safe_float(exec_summary.iloc[incremental_row, 3]) if incremental_row else 0.0
        incr_year2 = safe_float(exec_summary.iloc[incremental_row, 4]) if incremental_row else 0.0
        incr_year3 = safe_float(exec_summary.iloc[incremental_row, 5]) if incremental_row else 0.0
        
        pct_year1 = safe_float(exec_summary.iloc[pct_increase_row, 3]) if pct_increase_row else 0.0
        pct_year2 = safe_float(exec_summary.iloc[pct_increase_row, 4]) if pct_increase_row else 0.0
        pct_year3 = safe_float(exec_summary.iloc[pct_increase_row, 5]) if pct_increase_row else 0.0
        
        flex_year1 = safe_float(exec_summary.iloc[flexibility_final_row, 3]) if flexibility_final_row else 0.0
        flex_year2 = safe_float(exec_summary.iloc[flexibility_final_row, 4]) if flexibility_final_row else 0.0
        flex_year3 = safe_float(exec_summary.iloc[flexibility_final_row, 5]) if flexibility_final_row else 0.0
        
        # Load PowerPoint template
        print("4. Loading template...")
        prs = Presentation(template_pptx)
        
        # SLIDE 1
        print("5. Updating Slide 1...")
        slide1 = prs.slides[0]
        for shape in slide1.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "[Customer Name]" in run.text:
                            run.text = run.text.replace("[Customer Name]", customer_name)
        
        # SLIDE 9
        print("6. Updating Slide 9 (Font: 11pt)...")
        slide9 = prs.slides[8]
        
        for shape in slide9.shapes:
            if shape.shape_type == 14 and hasattr(shape, "text"):
                if "[Month Year]" in shape.text:
                    shape.text = "May 2026 compute, coverage & metrics"
        
        tables = [shape for shape in slide9.shapes if shape.shape_type == 19]
        
        if len(tables) >= 4:
            table1 = tables[0].table
            table1.rows[1].cells[1].text = f"{current_edp:.0%}"
            table1.rows[2].cells[1].text = f"{current_coverage:.1%}"
            table1.rows[3].cells[1].text = f"{current_flexibility:.2%}"
            table1.rows[4].cells[1].text = f"${current_savings:,.0f}"
            set_table_font_size(table1, 11)
            
            table2 = tables[1].table
            table2.rows[1].cells[1].text = f"{projected_edp:.0%}"
            table2.rows[2].cells[1].text = f"{projected_coverage:.1%}"
            table2.rows[3].cells[1].text = f"{projected_flexibility:.2%}"
            table2.rows[4].cells[1].text = f"${projected_savings:,.0f}"
            set_table_font_size(table2, 11)
            
            table3 = tables[2].table
            table3.rows[1].cells[1].text = f"{stable_pct:.2%}"
            table3.rows[2].cells[1].text = f"{fluctuations_pct:.2%}"
            table3.rows[3].cells[1].text = f"${stable_equivalent:,.2f}"
            set_table_font_size(table3, 11)
            
            table4 = tables[3].table
            table4.rows[0].cells[0].text = f"Customer: {customer_name}"
            set_table_font_size(table4, 11)
        
        # SLIDE 10
        print("7. Updating Slide 10 (Font: 8pt)...")
        slide10 = prs.slides[9]
        portfolio_tables = [shape for shape in slide10.shapes if shape.shape_type == 19]
        if portfolio_tables:
            for table_shape in portfolio_tables:
                set_table_font_size(table_shape.table, 8)
        
        # SLIDE 11
        print("8. Updating Slide 11 (Font: 11pt)...")
        slide11 = prs.slides[10]
        impact_tables = [shape for shape in slide11.shapes if shape.shape_type == 19]
        
        if impact_tables:
            impact_table = impact_tables[0].table
            
            impact_table.rows[1].cells[1].text = f"${current_year1:,.0f}"
            impact_table.rows[1].cells[2].text = f"${current_year2:,.0f}"
            impact_table.rows[1].cells[3].text = f"${current_year3:,.0f}"
            
            impact_table.rows[2].cells[1].text = f"${csa_year1:,.0f}"
            impact_table.rows[2].cells[2].text = f"${csa_year2:,.0f}"
            impact_table.rows[2].cells[3].text = f"${csa_year3:,.0f}"
            
            impact_table.rows[3].cells[1].text = f"${incr_year1:,.0f}"
            impact_table.rows[3].cells[2].text = f"${incr_year2:,.0f}"
            impact_table.rows[3].cells[3].text = f"${incr_year3:,.0f}"
            
            impact_table.rows[4].cells[1].text = f"{pct_year1:.2%}"
            impact_table.rows[4].cells[2].text = f"{pct_year2:.2%}"
            impact_table.rows[4].cells[3].text = f"{pct_year3:.2%}"
            
            impact_table.rows[5].cells[1].text = f"{flex_year1:.2%}"
            impact_table.rows[5].cells[2].text = f"{flex_year2:.2%}"
            impact_table.rows[5].cells[3].text = f"{flex_year3:.2%}"
            
            set_table_font_size(impact_table, 11)
        
        for shape in slide11.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if "[Customer]" in run.text:
                            run.text = run.text.replace("[Customer]", customer_name.split()[0])
        
        # Save
        print(f"9. Saving PowerPoint...")
        prs.save(output_pptx)
        
        print("\n" + "="*80)
        print("SUCCESS!")
        print("="*80)
        print(f"\nGenerated: {output_pptx}")
        print(f"\nSlide 9: Current Savings ${current_savings:,.0f}, Projected ${projected_savings:,.0f}")
        print(f"Slide 11: CSA Year 1 ${csa_year1:,.0f}, Year 2 ${csa_year2:,.0f}, Year 3 ${csa_year3:,.0f}")
        
    except Exception as e:
        print(f"\nERROR: {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)

if __name__ == "__main__":
    main()
